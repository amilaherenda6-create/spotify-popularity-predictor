#!/usr/bin/env python3
"""
generate_pptx.py — Spotify Popularity Predictor presentation
Dark navy / white / cyan color scheme, 13 slides.
Run: python reports/generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
BG     = RGBColor(0x1B, 0x2A, 0x4A)   # dark navy background
CARD   = RGBColor(0x24, 0x35, 0x58)   # card / panel
CARD2  = RGBColor(0x15, 0x22, 0x38)   # darker card
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xE8, 0xF4, 0xFD)   # very light blue body text
BLUE   = RGBColor(0xA8, 0xD4, 0xF5)   # secondary text
CYAN   = RGBColor(0x4F, 0xC3, 0xF7)   # heading accent
CYAN2  = RGBColor(0x00, 0xB4, 0xD8)   # second cyan
GREEN  = RGBColor(0x1D, 0xB9, 0x54)   # Spotify green
ORANGE = RGBColor(0xFF, 0x98, 0x00)   # warning / counter
MUTED  = RGBColor(0x55, 0x66, 0x77)   # dividers
TBEST  = RGBColor(0x07, 0x1F, 0x10)   # highlighted best-row background
THDR   = RGBColor(0x1E, 0x3A, 0x5F)   # table header row
TALT   = RGBColor(0x1D, 0x30, 0x54)   # alternating table row
TROW   = RGBColor(0x1B, 0x2A, 0x4A)   # normal table row

# Slide canvas (widescreen 16:9 = 10" × 5.625")
SW, SH = Inches(10), Inches(5.625)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
_blank = prs.slide_layouts[6]          # completely blank layout


# ── Low-level helpers ───────────────────────────────────────────────────────

def new_slide():
    """Blank slide with navy background + top/bottom cyan bars."""
    sl = prs.slides.add_slide(_blank)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    _rect(sl, 0, 0, SW, Inches(0.055), CYAN)
    _rect(sl, 0, SH - Inches(0.055), SW, Inches(0.055), CYAN2)
    return sl


def _rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _tx(slide, text, l, t, w, h, sz=12, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True):
    color = color or WHITE
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        rn = p.add_run()
        rn.text = line
        rn.font.name = "Calibri"
        rn.font.size = Pt(sz)
        rn.font.bold = bold
        rn.font.italic = italic
        rn.font.color.rgb = color
    return box


def heading(slide, text, y=Inches(0.2)):
    _tx(slide, text, Inches(0.5), y, Inches(9), Inches(0.5),
        sz=26, bold=True, color=CYAN)
    _rect(slide, Inches(0.5), y + Inches(0.5), Inches(9), Inches(0.035), CYAN)


def tbl(slide, hdrs, rows, l, t, w, h, best=None, hsz=10, rsz=9.5):
    """Add a styled table. best=row-index (0-based) highlights that data row."""
    nr, nc = len(rows) + 1, len(hdrs)
    tshp = slide.shapes.add_table(nr, nc, l, t, w, h)
    tb = tshp.table

    for c, hdr in enumerate(hdrs):
        cell = tb.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = THDR
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        rn = p.add_run()
        rn.text = hdr
        rn.font.bold = True
        rn.font.size = Pt(hsz)
        rn.font.color.rgb = CYAN
        rn.font.name = "Calibri"

    for r, row in enumerate(rows):
        is_alt  = r % 2 == 1
        is_best = best is not None and r == best
        for c, val in enumerate(row):
            cell = tb.cell(r + 1, c)
            cell.fill.solid()
            if is_best:
                cell.fill.fore_color.rgb = TBEST
            elif is_alt:
                cell.fill.fore_color.rgb = TALT
            else:
                cell.fill.fore_color.rgb = TROW
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            rn = p.add_run()
            rn.text = str(val)
            rn.font.size = Pt(rsz)
            rn.font.bold = is_best
            rn.font.color.rgb = GREEN if is_best else LIGHT
            rn.font.name = "Calibri"
    return tb


def card(slide, l, t, w, h, accent=CYAN, bg=None):
    bg = bg or CARD
    _rect(slide, l, t, w, h, bg)
    _rect(slide, l, t, Inches(0.06), h, accent)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
_rect(sl, 0, 0, SW, Inches(0.055), CYAN)

_tx(sl, "Spotify Popularity Predictor",
    Inches(0.5), Inches(0.65), Inches(9), Inches(1.1),
    sz=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

_rect(sl, Inches(3.5), Inches(1.6), Inches(3), Inches(0.05), GREEN)

_tx(sl, "Can audio features predict a song's success?",
    Inches(0.5), Inches(1.72), Inches(9), Inches(0.45),
    sz=15, color=BLUE, align=PP_ALIGN.CENTER)

_rect(sl, Inches(1.3), Inches(2.35), Inches(7.4), Inches(2.85), CARD)

info = [
    ("Student",     "Herenda Amila"),
    ("Course",      "Modelling in Advanced Data Analytics"),
    ("Professors",  "Uroš Godnov · Aleš Gorišek"),
    ("Programme",   "Int'l Full-Time Master's in Business & Organisation"),
    ("Date",        "03.07.2026"),
    ("University",  "FELU — Faculty of Economics Ljubljana"),
]
for i, (lbl, val) in enumerate(info):
    y = Inches(2.48) + Inches(i * 0.4)
    _tx(sl, lbl + ":", Inches(1.5), y, Inches(1.7), Inches(0.37),
        sz=10, bold=True, color=CYAN)
    _tx(sl, val, Inches(3.3), y, Inches(5.1), Inches(0.37),
        sz=10, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "The Problem")

_rect(sl, Inches(0.5), Inches(0.92), Inches(9), Inches(1.0), CARD2)
_rect(sl, Inches(0.5), Inches(0.92), Inches(0.06), Inches(1.0), CYAN)
_tx(sl, "“Every year, millions of songs are uploaded to Spotify.\n"
        "Most are never heard. A handful become global hits.”",
    Inches(0.7), Inches(0.98), Inches(8.6), Inches(0.88),
    sz=13, italic=True, color=BLUE)

_tx(sl, "The question every artist, producer, and label manager asks:",
    Inches(0.5), Inches(2.07), Inches(9), Inches(0.3),
    sz=11, bold=True, color=LIGHT)

card(sl, Inches(0.5), Inches(2.45), Inches(4.3), Inches(0.9), accent=GREEN,
     bg=RGBColor(0x07, 0x18, 0x10))
_tx(sl, "\U0001f3af  Will this song be a hit?",
    Inches(0.7), Inches(2.52), Inches(4.0), Inches(0.38),
    sz=13, bold=True, color=WHITE)
_tx(sl, "Yes / No prediction",
    Inches(0.7), Inches(2.88), Inches(4.0), Inches(0.3),
    sz=10, color=BLUE)

card(sl, Inches(5.2), Inches(2.45), Inches(4.3), Inches(0.9), accent=CYAN,
     bg=RGBColor(0x06, 0x10, 0x1c))
_tx(sl, "\U0001f4ca  How popular will it be?",
    Inches(5.4), Inches(2.52), Inches(4.0), Inches(0.38),
    sz=13, bold=True, color=WHITE)
_tx(sl, "Score 0 – 100",
    Inches(5.4), Inches(2.88), Inches(4.0), Inches(0.3),
    sz=10, color=BLUE)

_tx(sl, "Our approach: Analyse only the audio characteristics of a song —\n"
        "no artist name, no label, no marketing budget.\n"
        "If the sound contains the secret to success, our tool will find it.",
    Inches(0.5), Inches(3.52), Inches(9), Inches(0.9),
    sz=11, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Data
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "The Data")

_tx(sl, "Source: Spotify Tracks Dataset (Kaggle 2022)  ·  114,000 tracks  ·  21 audio features  ·  114 genres",
    Inches(0.5), Inches(0.88), Inches(9), Inches(0.3),
    sz=10, color=BLUE)

stats = [
    ("114,000", "Spotify songs"),
    ("114",     "music genres"),
    ("21",      "features per song"),
    ("0–100", "popularity target"),
]
for i, (num, label) in enumerate(stats):
    y = Inches(1.32) + Inches(i * 0.93)
    _rect(sl, Inches(0.5), y, Inches(3.8), Inches(0.82), CARD)
    _tx(sl, num, Inches(0.68), y + Inches(0.04), Inches(2.0), Inches(0.46),
        sz=26, bold=True, color=CYAN)
    _tx(sl, label, Inches(0.68), y + Inches(0.46), Inches(3.4), Inches(0.3),
        sz=10, color=BLUE)

tbl(sl,
    ["Feature", "What it measures"],
    [
        ["Danceability", "Suitability for dancing"],
        ["Energy",       "Intensity and activity"],
        ["Loudness",     "Overall loudness (dB)"],
        ["Valence",      "Positiveness / mood"],
        ["Tempo",        "Speed in beats per minute"],
        ["Acousticness", "Live / organic vs electronic"],
    ],
    l=Inches(4.55), t=Inches(1.28), w=Inches(5.2), h=Inches(3.9),
    hsz=10, rsz=9.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Our Approach
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "Our Approach")
_tx(sl, "A practical end-to-end prediction workflow:",
    Inches(0.5), Inches(0.85), Inches(9), Inches(0.3),
    sz=12, color=LIGHT)

steps = [
    ("1", "Clean the dataset, remove duplicates, handle missing values"),
    ("2", "Train classification and regression models independently"),
    ("3", "Compare each model against a simple baseline (random guessing)"),
    ("4", "Deploy a live app for instant song scoring"),
]
for i, (num, desc) in enumerate(steps):
    y = Inches(1.25) + Inches(i * 0.7)
    _rect(sl, Inches(0.5), y, Inches(0.55), Inches(0.55), CYAN)
    _tx(sl, num, Inches(0.5), y + Inches(0.06), Inches(0.55), Inches(0.42),
        sz=18, bold=True, color=BG, align=PP_ALIGN.CENTER)
    _tx(sl, desc, Inches(1.18), y + Inches(0.1), Inches(8.4), Inches(0.4),
        sz=12, color=LIGHT)

_tx(sl, "We define a HIT as popularity score ≥ 70  (top 5% of all songs on Spotify)",
    Inches(0.5), Inches(4.15), Inches(9), Inches(0.3),
    sz=11, italic=True, color=BLUE)

# Two goal boxes
card(sl, Inches(0.5), Inches(4.55), Inches(4.3), Inches(0.7),
     accent=GREEN, bg=RGBColor(0x07, 0x18, 0x10))
_tx(sl, "\U0001f3af  CLASSIFICATION:  Hit or not?  (Yes / No)",
    Inches(0.7), Inches(4.63), Inches(4.0), Inches(0.5),
    sz=11, bold=True, color=WHITE)

card(sl, Inches(5.2), Inches(4.55), Inches(4.3), Inches(0.7),
     accent=CYAN, bg=RGBColor(0x06, 0x10, 0x1c))
_tx(sl, "\U0001f4ca  REGRESSION:  Exact score 0 – 100",
    Inches(5.4), Inches(4.63), Inches(4.0), Inches(0.5),
    sz=11, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Classification Goal
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "Classification Goal: Can this song be popular?")

card(sl, Inches(0.5), Inches(0.95), Inches(4.3), Inches(1.25),
     accent=GREEN, bg=RGBColor(0x07, 0x18, 0x10))
_tx(sl, "Popular (YES)", Inches(0.7), Inches(1.02), Inches(3.8), Inches(0.38),
    sz=16, bold=True, color=GREEN)
_tx(sl, "Popularity score ≥ 70", Inches(0.7), Inches(1.38), Inches(3.8), Inches(0.3),
    sz=12, color=LIGHT)
_tx(sl, "Top 5% on Spotify", Inches(0.7), Inches(1.66), Inches(3.8), Inches(0.3),
    sz=11, color=BLUE)

card(sl, Inches(5.2), Inches(0.95), Inches(4.3), Inches(1.25),
     accent=ORANGE, bg=RGBColor(0x1A, 0x08, 0x02))
_tx(sl, "Not Popular (NO)", Inches(5.4), Inches(1.02), Inches(3.8), Inches(0.38),
    sz=16, bold=True, color=ORANGE)
_tx(sl, "Popularity score < 70", Inches(5.4), Inches(1.38), Inches(3.8), Inches(0.3),
    sz=12, color=LIGHT)
_tx(sl, "95% of all songs uploaded", Inches(5.4), Inches(1.66), Inches(3.8), Inches(0.3),
    sz=11, color=BLUE)

_tx(sl, "Why does class imbalance matter?",
    Inches(0.5), Inches(2.42), Inches(9), Inches(0.3),
    sz=12, bold=True, color=LIGHT)

points = [
    "Only 1 in 20 songs is a hit — so the dataset is heavily skewed",
    "A naive model that says NO to everything gets 95% accuracy — but finds zero hits",
    "We use F1-score and AUC instead of accuracy as our honest success metrics",
    "XGBoost uses scale_pos_weight to correct the imbalance during training",
]
for i, pt in enumerate(points):
    y = Inches(2.82) + Inches(i * 0.52)
    _tx(sl, "▸  " + pt, Inches(0.5), y, Inches(9.0), Inches(0.46),
        sz=11, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Regression Goal
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "Regression Goal: How popular will the song score?")

_tx(sl, "Instead of YES/NO, we predict the exact Spotify popularity score (0 to 100).",
    Inches(0.5), Inches(0.9), Inches(9), Inches(0.3),
    sz=12, color=LIGHT)

ranges = [
    ("0 – 20",  MUTED,    "Unknown / brand-new songs"),
    ("20 – 40", BLUE,     "Most songs — moderate play"),
    ("40 – 60", CYAN,     "Gaining traction, some fans"),
    ("60 – 70", GREEN,    "Strong performer"),
    ("70 – 100", RGBColor(0xFF, 0xD7, 0x00), "HIT: top 5% on the platform"),
]
for i, (rng, col, desc) in enumerate(ranges):
    y = Inches(1.35) + Inches(i * 0.68)
    _rect(sl, Inches(0.5), y + Inches(0.1), Inches(1.5), Inches(0.42), col)
    _tx(sl, rng, Inches(0.5), y + Inches(0.1), Inches(1.5), Inches(0.42),
        sz=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
    _tx(sl, desc, Inches(2.2), y + Inches(0.16), Inches(7.3), Inches(0.36),
        sz=11, color=LIGHT)

_tx(sl, "Key metric: RMSE (average error in popularity points)\n"
        "XGBoost achieves RMSE = 16.7 — on average 17 points from the real score.",
    Inches(0.5), Inches(4.78), Inches(9), Inches(0.5),
    sz=11, italic=True, color=BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Data Insights
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "What We Found in the Data")

findings = [
    ("1 in 20",      CYAN,   "Songs reach the hit threshold\n5,465 popular out of 106,907 tracks"),
    ("20 – 40", BLUE,   "Most songs cluster in this range\nMany new or obscure tracks score near 0"),
    ("Genre #1",     GREEN,  "Genre is the strongest single predictor\nPop, dance, hip-hop score highest"),
    ("Sound wins",   ORANGE, "Loud + energetic + danceable = better\nAcoustic + instrumental = lower score"),
]
positions = [
    (Inches(0.5), Inches(4.6)),
    (Inches(5.2), Inches(4.6)),
    (Inches(0.5), Inches(4.6)),
    (Inches(5.2), Inches(4.6)),
]
tops = [Inches(1.1), Inches(1.1), Inches(3.05), Inches(3.05)]
lefts = [Inches(0.5), Inches(5.2), Inches(0.5), Inches(5.2)]

for i, (num, col, desc) in enumerate(findings):
    l = lefts[i]
    t = tops[i]
    _rect(sl, l, t, Inches(4.3), Inches(1.7), CARD)
    _rect(sl, l, t, Inches(0.06), Inches(1.7), col)
    _tx(sl, num, l + Inches(0.16), t + Inches(0.08), Inches(3.8), Inches(0.52),
        sz=22, bold=True, color=col)
    _tx(sl, desc, l + Inches(0.16), t + Inches(0.62), Inches(4.0), Inches(0.96),
        sz=10, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Performance
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "How Well Does It Work?")

_tx(sl, "YES / NO — Can we identify a hit?",
    Inches(0.5), Inches(0.88), Inches(9), Inches(0.3),
    sz=11, bold=True, color=LIGHT)
tbl(sl,
    ["Approach", "Hit detection (Recall)", "Reliability (AUC)"],
    [
        ["Random guessing",      "0% — misses every hit", "0.50"],
        ["Logistic Regression",  "Some hits found",            "0.71"],
        ["Decision Tree",        "Similar to above",           "0.71"],
        ["XGBoost (best) ✅", "76 in 100 hits found",     "0.92"],
    ],
    l=Inches(0.5), t=Inches(1.2), w=Inches(9), h=Inches(1.65),
    best=3, hsz=10, rsz=9.5)

_tx(sl, "Exact score — how close are our guesses?",
    Inches(0.5), Inches(3.0), Inches(9), Inches(0.3),
    sz=11, bold=True, color=LIGHT)
tbl(sl,
    ["Approach", "Average error (RMSE)", "vs baseline"],
    [
        ["Always predict average",   "±22 points", "—"],
        ["Ridge Regression",          "±21 points", "4% better"],
        ["Decision Tree",             "±21 points", "9% better"],
        ["XGBoost (best) ✅",     "±17 points", "38% better"],
    ],
    l=Inches(0.5), t=Inches(3.32), w=Inches(9), h=Inches(1.65),
    best=3, hsz=10, rsz=9.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — What Matters Most
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "What Matters Most?")
_tx(sl, "Audio signals that most influence popularity score:",
    Inches(0.5), Inches(0.85), Inches(9), Inches(0.3),
    sz=11, color=LIGHT)

# Boosts column
card(sl, Inches(0.5), Inches(1.22), Inches(4.3), Inches(2.9),
     accent=GREEN, bg=RGBColor(0x07, 0x18, 0x10))
_tx(sl, "Boosts popularity ↑",
    Inches(0.7), Inches(1.28), Inches(4.0), Inches(0.38),
    sz=13, bold=True, color=GREEN)
boosts = [
    "\U0001f3b5  Genre — pop, dance, hip-hop lead",
    "\U0001f50a  Loudness — louder = more streams",
    "\U0001f483  Danceability — playlists reward it",
    "⚡  Energy — intense tracks perform better",
]
for i, b in enumerate(boosts):
    _tx(sl, b, Inches(0.7), Inches(1.74) + Inches(i * 0.5),
        Inches(4.0), Inches(0.44), sz=10.5, color=LIGHT)

# Reduces column
card(sl, Inches(5.2), Inches(1.22), Inches(4.3), Inches(2.9),
     accent=CYAN, bg=RGBColor(0x06, 0x10, 0x20))
_tx(sl, "Reduces popularity ↓",
    Inches(5.4), Inches(1.28), Inches(4.0), Inches(0.38),
    sz=13, bold=True, color=CYAN)
reduces = [
    "\U0001f3b9  Instrumentalness — no vocals = lower",
    "\U0001f3b8  Acousticness — electronic beats win",
    "\U0001f634  Duration — very long tracks score lower",
    "\U0001f399️  Speechiness — too much talking hurts",
]
for i, r in enumerate(reduces):
    _tx(sl, r, Inches(5.4), Inches(1.74) + Inches(i * 0.5),
        Inches(4.0), Inches(0.44), sz=10.5, color=LIGHT)

_rect(sl, Inches(0.5), Inches(4.28), Inches(9), Inches(0.6), CARD)
_tx(sl, "\U0001f4a1  New combined feature: “Club score” (danceability × energy) "
        "improved predictions beyond any single measurement.",
    Inches(0.65), Inches(4.34), Inches(8.6), Inches(0.46),
    sz=10, color=BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Deployment
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "Try It Live")

_tx(sl, "\U0001f310  Open the app now:",
    Inches(0.5), Inches(0.88), Inches(9), Inches(0.32),
    sz=12, bold=True, color=LIGHT)
_tx(sl, "spotify-popularity-predictor-dkb249nsp7axa3hkxdpkht.streamlit.app",
    Inches(0.5), Inches(1.18), Inches(9), Inches(0.38),
    sz=13, bold=True, color=CYAN)
_rect(sl, Inches(0.5), Inches(1.54), Inches(9), Inches(0.035), CYAN2)

# Popular card
card(sl, Inches(0.5), Inches(1.7), Inches(4.3), Inches(2.85),
     accent=GREEN, bg=RGBColor(0x07, 0x1A, 0x10))
_tx(sl, "HIGH-POPULARITY PREDICTION", Inches(0.7), Inches(1.78), Inches(4.0), Inches(0.3),
    sz=9, bold=True, color=BLUE)
_tx(sl, "Pop song: high danceability + energy", Inches(0.7), Inches(2.06), Inches(4.0), Inches(0.28),
    sz=9, italic=True, color=BLUE)
_tx(sl, "✅  POPULAR", Inches(0.7), Inches(2.38), Inches(4.0), Inches(0.5),
    sz=22, bold=True, color=GREEN)
_tx(sl, "Confidence: 92.6%\nPredicted score: 53.6 / 100",
    Inches(0.7), Inches(2.9), Inches(4.0), Inches(0.55),
    sz=11, color=LIGHT)

# Not popular card
card(sl, Inches(5.2), Inches(1.7), Inches(4.3), Inches(2.85),
     accent=ORANGE, bg=RGBColor(0x1A, 0x0A, 0x02))
_tx(sl, "LOW-POPULARITY PREDICTION", Inches(5.4), Inches(1.78), Inches(4.0), Inches(0.3),
    sz=9, bold=True, color=BLUE)
_tx(sl, "Classical instrumental: quiet + acoustic", Inches(5.4), Inches(2.06), Inches(4.0), Inches(0.28),
    sz=9, italic=True, color=BLUE)
_tx(sl, "⚠️  NOT POPULAR", Inches(5.4), Inches(2.38), Inches(4.0), Inches(0.5),
    sz=22, bold=True, color=ORANGE)
_tx(sl, "Confidence: 0.07%\nPredicted score: 20.8 / 100",
    Inches(5.4), Inches(2.9), Inches(4.0), Inches(0.55),
    sz=11, color=LIGHT)

_tx(sl, "Adjust any audio slider — get an instant prediction",
    Inches(0.5), Inches(4.7), Inches(9), Inches(0.3),
    sz=10, italic=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Limitations
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "Honest Limitations")

card(sl, Inches(0.5), Inches(0.95), Inches(4.4), Inches(2.55),
     accent=GREEN, bg=RGBColor(0x07, 0x18, 0x10))
_tx(sl, "What it does well  ✅",
    Inches(0.7), Inches(1.0), Inches(4.0), Inches(0.38),
    sz=12, bold=True, color=GREEN)
for i, d in enumerate([
    "Identifies popular vs unpopular reliably",
    "Finds 76 out of every 100 genuine hits",
    "Good first filter for large catalogues",
    "Confirms genre + energy + danceability matter",
]):
    _tx(sl, "✓  " + d, Inches(0.7), Inches(1.45) + Inches(i * 0.48),
        Inches(4.0), Inches(0.42), sz=10, color=LIGHT)

card(sl, Inches(5.1), Inches(0.95), Inches(4.4), Inches(2.55),
     accent=ORANGE, bg=RGBColor(0x1A, 0x08, 0x08))
_tx(sl, "What it cannot do  ❌",
    Inches(5.3), Inches(1.0), Inches(4.0), Inches(0.38),
    sz=12, bold=True, color=ORANGE)
for i, c in enumerate([
    "Predict viral social momentum",
    "Evaluate artistic quality",
    "Work well for niche genres",
    "Replace expert human judgment",
]):
    _tx(sl, "✗  " + c, Inches(5.3), Inches(1.45) + Inches(i * 0.48),
        Inches(4.0), Inches(0.42), sz=10, color=LIGHT)

_rect(sl, Inches(0.5), Inches(3.68), Inches(9), Inches(1.45), CARD)
_tx(sl, "38%", Inches(0.65), Inches(3.76), Inches(1.7), Inches(0.82),
    sz=44, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
_tx(sl, "of popularity explained\nby audio features alone",
    Inches(2.3), Inches(3.84), Inches(3.0), Inches(0.72),
    sz=11, color=LIGHT)
_rect(sl, Inches(5.1), Inches(3.8), Inches(0.04), Inches(1.15), MUTED)
_tx(sl, "62%", Inches(5.3), Inches(3.76), Inches(1.7), Inches(0.82),
    sz=44, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
_tx(sl, "artist fame, marketing,\ntiming & everything else",
    Inches(7.0), Inches(3.84), Inches(2.6), Inches(0.72),
    sz=11, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Recommendations
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
heading(sl, "Recommendations")

card(sl, Inches(0.5), Inches(0.95), Inches(4.3), Inches(3.95),
     accent=GREEN, bg=RGBColor(0x07, 0x18, 0x10))
_tx(sl, "For artists & producers",
    Inches(0.7), Inches(1.0), Inches(4.0), Inches(0.38),
    sz=13, bold=True, color=GREEN)
for i, r in enumerate([
    "Choose genre intentionally",
    "Aim for high energy and danceability",
    "Keep tracks under 4 minutes",
    "Produce loudly (within streaming norms)",
]):
    _tx(sl, "▸  " + r, Inches(0.7), Inches(1.48) + Inches(i * 0.6),
        Inches(4.0), Inches(0.52), sz=11, color=LIGHT)
_tx(sl, "These are data patterns — not rules.\nOriginality still matters.",
    Inches(0.7), Inches(3.95), Inches(4.0), Inches(0.7),
    sz=10, italic=True, color=BLUE)

card(sl, Inches(5.2), Inches(0.95), Inches(4.3), Inches(3.95),
     accent=CYAN, bg=RGBColor(0x06, 0x0e, 0x1c))
_tx(sl, "For labels & A&R teams",
    Inches(5.4), Inches(1.0), Inches(4.0), Inches(0.38),
    sz=13, bold=True, color=CYAN)
for i, r in enumerate([
    "Use as a first-pass filter on demos",
    "Surface new tracks before listen data exists",
    "Do not use as the final gatekeeper",
    "Human taste and expertise remain essential",
]):
    _tx(sl, "▸  " + r, Inches(5.4), Inches(1.48) + Inches(i * 0.6),
        Inches(4.0), Inches(0.52), sz=11, color=LIGHT)
_tx(sl, "Treat prediction as guidance,\nnot the final decision.",
    Inches(5.4), Inches(3.95), Inches(4.0), Inches(0.7),
    sz=10, italic=True, color=BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
_rect(sl, Inches(0.5), Inches(0.18), Inches(9), Inches(5.05), CARD)
_rect(sl, Inches(0.5), Inches(0.18), Inches(9), Inches(0.06), GREEN)

_tx(sl, "Thank You",
    Inches(0.5), Inches(0.42), Inches(9), Inches(1.1),
    sz=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

_rect(sl, Inches(3.5), Inches(1.35), Inches(3.0), Inches(0.05), GREEN)

_tx(sl, "Herenda Amila",
    Inches(0.5), Inches(1.5), Inches(9), Inches(0.45),
    sz=18, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

for i, line in enumerate([
    "Modelling in Advanced Data Analytics",
    "Faculty of Economics Ljubljana (FELU)",
    "Int'l Full-Time Master's Programme in Business and Organisation",
    "03.07.2026",
]):
    _tx(sl, line, Inches(0.5), Inches(2.02) + Inches(i * 0.38),
        Inches(9), Inches(0.35), sz=10.5, color=BLUE, align=PP_ALIGN.CENTER)

_rect(sl, Inches(2.3), Inches(3.62), Inches(5.4), Inches(0.58), BG)
_tx(sl, "Professors: Uroš Godnov  ·  Aleš Gorišek",
    Inches(2.3), Inches(3.68), Inches(5.4), Inches(0.42),
    sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

_tx(sl, "\U0001f310  Try the app live:",
    Inches(0.5), Inches(4.34), Inches(9), Inches(0.28),
    sz=10, color=LIGHT, align=PP_ALIGN.CENTER)
_tx(sl, "spotify-popularity-predictor-dkb249nsp7axa3hkxdpkht.streamlit.app",
    Inches(0.5), Inches(4.6), Inches(9), Inches(0.35),
    sz=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(__file__), "spotify_predictor.pptx")
prs.save(out)
print(f"Saved {prs.slides.__len__()} slides -> {out}")
